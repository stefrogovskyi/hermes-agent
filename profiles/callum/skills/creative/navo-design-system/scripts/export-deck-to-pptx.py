#!/usr/bin/env python3
"""
Export web slide decks to pixel-perfect, non-editable PowerPoint (.pptx).
Usage: python scripts/export-deck-to-pptx.py <deck_url> <output_pptx> [--slides 14] [--mode laptop|fullbleed]

Default mode 'laptop':
  Captures standard 13-inch laptop viewport (1440x900 Retina 2x) directly from the web view.
  Preserves exact proportions, large readable fonts, and comfortable card spacing without CSS overrides.
  Sets presentation slide size to 14.4" x 9.0" (16:10).

Mode 'fullbleed':
  Injects custom presentation CSS for 1920x1080 16:9 widescreen, ensuring strict slide isolation (.slide.active ONLY).
"""

import argparse
import os
import sys
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches

def export_deck(url: str, output_path: str, total_slides: int, mode: str = "laptop"):
    tmp_dir = "/tmp/deck_export_frames"
    os.makedirs(tmp_dir, exist_ok=True)

    if mode == "laptop":
        vp_width, vp_height = 1440, 900
        slide_w_in, slide_h_in = 14.4, 9.0
        print(f"[1/2] Capturing {total_slides} slides at standard 13-inch laptop resolution ({vp_width}x{vp_height} Retina 2x)...")
    else:
        vp_width, vp_height = 1920, 1080
        slide_w_in, slide_h_in = 13.333333, 7.5
        print(f"[1/2] Capturing {total_slides} slides in Full-Bleed 16:9 ({vp_width}x{vp_height} Retina 2x)...")

    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page(
            viewport={"width": vp_width, "height": vp_height},
            device_scale_factor=2
        )
        page.goto(url)
        page.wait_for_timeout(2000)

        if mode == "fullbleed":
            # Guard against slide pileup: .slide must be display: none, only .slide.active display: flex!
            page.evaluate("""() => {
                const style = document.createElement("style");
                style.innerHTML = `
                    .deck-header, .progress-bar-container, .mobile-controls, .nav-btn, .theme-toggle-btn {
                        display: none !important;
                    }
                    .deck-stage, body, html {
                        padding: 0 !important; margin: 0 !important;
                        width: 1920px !important; height: 1080px !important;
                        overflow: hidden !important; display: block !important;
                    }
                    .slide {
                        display: none !important;
                        width: 1920px !important; max-width: 1920px !important;
                        height: 1080px !important; max-height: 1080px !important;
                        border-radius: 0 !important; border: none !important; box-shadow: none !important;
                        margin: 0 !important; padding: 48px 72px 40px 72px !important;
                        box-sizing: border-box !important;
                    }
                    .slide.active {
                        display: flex !important; flex-direction: column !important; justify-content: space-between !important;
                    }
                    .slide-title { font-size: 3.1rem !important; margin-bottom: 14px !important; line-height: 1.15 !important; }
                    .slide-subhead { font-size: 1.2rem !important; max-width: 1400px !important; margin-bottom: 24px !important; line-height: 1.5 !important; }
                    .cards-grid-3 { gap: 24px !important; }
                    .cards-grid-2 { gap: 28px !important; }
                    .cards-grid-5 { gap: 18px !important; }
                    .deck-card { padding: 26px 28px !important; border-radius: 20px !important; }
                    .card-heading { font-size: 1.35rem !important; }
                    .card-text { font-size: 1.02rem !important; line-height: 1.55 !important; }
                    .screenshot-showcase-container { margin-top: 12px !important; border-radius: 20px !important; }
                    .screenshot-showcase-img { max-height: 560px !important; }
                `;
                document.head.appendChild(style);

                const slides = document.querySelectorAll(".slide");
                slides.forEach(slide => {
                    const headerRow = slide.querySelector(".slide-header-row");
                    if (headerRow && !slide.querySelector(".injected-brand-logo")) {
                        const brandSvg = document.querySelector(".navo-brand-logo");
                        if (brandSvg) {
                            const cloneSvg = brandSvg.cloneNode(true);
                            cloneSvg.classList.add("injected-brand-logo");
                            cloneSvg.style.height = "26px";
                            cloneSvg.style.width = "auto";
                            const leftDiv = document.createElement("div");
                            leftDiv.style.display = "flex";
                            leftDiv.style.alignItems = "center";
                            leftDiv.style.gap = "14px";
                            const eyebrow = headerRow.querySelector(".slide-eyebrow");
                            if (eyebrow) {
                                headerRow.insertBefore(leftDiv, eyebrow);
                                leftDiv.appendChild(cloneSvg);
                                const sep = document.createElement("span");
                                sep.textContent = "|";
                                sep.style.color = "rgba(255,255,255,0.2)";
                                sep.style.fontSize = "1.2rem";
                                sep.style.fontWeight = "300";
                                leftDiv.appendChild(sep);
                                leftDiv.appendChild(eyebrow);
                            }
                        }
                    }
                });
            }""")

        for i in range(1, total_slides + 1):
            page.evaluate(f"if (typeof updateSlide === 'function') updateSlide({i});")
            page.wait_for_timeout(600)
            frame_path = os.path.join(tmp_dir, f"slide_{i:02d}.png")
            page.screenshot(path=frame_path)
            print(f"  ✓ Slide {i:02d} rendered")

        browser.close()

    print(f"[2/2] Assembling non-editable PPTX at {output_path}...")
    prs = Presentation()
    prs.slide_width = Inches(slide_w_in)
    prs.slide_height = Inches(slide_h_in)
    blank_layout = prs.slide_layouts[6]

    for i in range(1, total_slides + 1):
        frame_path = os.path.join(tmp_dir, f"slide_{i:02d}.png")
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            frame_path,
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height
        )

    prs.save(output_path)
    size_mb = os.path.getsize(output_path) / (1024 * 1024)
    print(f"Done: {output_path} ({size_mb:.2f} MB)")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Export web deck to pixel-perfect non-editable PPTX")
    parser.add_argument("url", help="URL of the web slide deck")
    parser.add_argument("output", help="Path to save the .pptx file")
    parser.add_argument("--slides", type=int, default=14, help="Total number of slides")
    parser.add_argument("--mode", choices=["laptop", "fullbleed"], default="laptop", help="Capture mode")
    args = parser.parse_args()

    export_deck(args.url, args.output, args.slides, args.mode)
