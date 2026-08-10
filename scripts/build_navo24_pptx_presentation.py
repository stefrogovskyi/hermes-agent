# -*- coding: utf-8 -*-
"""
build_navo24_pptx_presentation.py — Генератор нативной PowerPoint презентации (.pptx)
в дизайн-системе Navo24 (темная тема, высоконтрастный дизайн, 10 слайдов).
"""

import os, pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

REPORTS_DIR = r"C:\Users\Stefan\AppData\Local\hermes\reports"
os.makedirs(REPORTS_DIR, exist_ok=True)

pptx_path = os.path.join(REPORTS_DIR, "navo24_team_presentation.pptx")

prs = pptx.Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette (Navo24 Design System)
BG_COLOR = RGBColor(11, 15, 25)      # #0B0F19 Dark Void
CARD_BG = RGBColor(17, 24, 39)       # #111827 Dark Navy Card
SKY_BLUE = RGBColor(2, 132, 199)      # #0284C7 Accent
CYAN_GLOW = RGBColor(56, 189, 248)    # #38BDF8 Glow
TEXT_WHITE = RGBColor(248, 250, 252) # #F8FAFC
TEXT_MUTED = RGBColor(148, 163, 184) # #94A3B8
EMERALD = RGBColor(16, 185, 129)     # #10B981 Success
AMBER = RGBColor(245, 158, 11)       # #F59E0B Warning/Spartan

blank_layout = prs.slide_layouts[6]

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_header(slide, tag_text, title_text, subtitle_text=""):
    set_slide_background(slide)
    
    # Header Box
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.8))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p_tag = tf.paragraphs[0]
    p_tag.text = tag_text.upper()
    p_tag.font.size = Pt(12)
    p_tag.font.bold = True
    p_tag.font.color.rgb = CYAN_GLOW
    
    p_title = tf.add_paragraph()
    p_title.text = title_text
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_WHITE
    
    if subtitle_text:
        p_sub = tf.add_paragraph()
        p_sub.text = subtitle_text
        p_sub.font.size = Pt(14)
        p_sub.font.color.rgb = TEXT_MUTED

# Slide 1: Title Slide
slide1 = prs.slides.add_slide(blank_layout)
set_slide_background(slide1)

tb1 = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(3.5))
tf1 = tb1.text_frame
tf1.word_wrap = True

p_tag = tf1.paragraphs[0]
p_tag.text = "EXECUTIVE STRATEGY BRIEFING"
p_tag.font.size = Pt(14)
p_tag.font.bold = True
p_tag.font.color.rgb = CYAN_GLOW

p_title = tf1.add_paragraph()
p_title.text = "The Next Chapter:\nOur Path to Break Free and Win"
p_title.font.size = Pt(40)
p_title.font.bold = True
p_title.font.color.rgb = TEXT_WHITE

p_sub = tf1.add_paragraph()
p_sub.text = "A candid message on transition, ownership, and our limitless future."
p_sub.font.size = Pt(18)
p_sub.font.color.rgb = CYAN_GLOW

# Slide 2: Core Mindset
slide2 = prs.slides.add_slide(blank_layout)
add_header(slide2, "Slide 02 / Core Mindset", "Obsession with Customer Success", "Shift Your Focus: From Transactions to Lifelong Relationships")

tb2 = slide2.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(4.0))
tf2 = tb2.text_frame
tf2.word_wrap = True

p = tf2.paragraphs[0]
p.text = "• Customers for Life: We aren't just selling tools; we are obsessively committed to solving customer problems."
p.font.size = Pt(18)
p.font.color.rgb = TEXT_WHITE

p2 = tf2.add_paragraph()
p2.text = "• Service as Primary Growth Driver: Exceptional, proactive service is our primary engine for long-term expansion."
p2.font.size = Pt(18)
p2.font.color.rgb = TEXT_MUTED

# Slide 3: Earning Power
slide3 = prs.slides.add_slide(blank_layout)
add_header(slide3, "Slide 03 / Potential & Earning Power", "$10,000+ Monthly Earnings", "A Realistic Target Built on Ownership, Not a Dream")

tb3 = slide3.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(4.0))
tf3 = tb3.text_frame
tf3.word_wrap = True

p = tf3.paragraphs[0]
p.text = "• Absolute Reality: We are 100% independent, AI & tech-empowered, and completely free from corporate red tape."
p.font.size = Pt(18)
p.font.color.rgb = EMERALD

p2 = tf3.add_paragraph()
p2.text = "• Step Out of Comfort: Leave the comfortable routine behind—your results now depend entirely on your drive."
p2.font.size = Pt(18)
p2.font.color.rgb = TEXT_WHITE

p3 = tf3.add_paragraph()
p3.text = "• Unshakeable Resilience: Early rejections will happen, but every setback is just data to improve."
p3.font.size = Pt(18)
p3.font.color.rgb = TEXT_MUTED

# Slide 4: Us vs Corporate Giant
slide4 = prs.slides.add_slide(blank_layout)
add_header(slide4, "Slide 04 / Competitive Advantage", "Why We Break Through the Ceiling", "Speed, Agility, and Total Control vs. Corporate Giants")

tb4 = slide4.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.7), Inches(4.2))
tf4 = tb4.text_frame
tf4.word_wrap = True

p = tf4.paragraphs[0]
p.text = "1. Real-time product updates (deployed instantly vs months of committee reviews)"
p.font.size = Pt(16)
p.font.color.rgb = TEXT_WHITE

p2 = tf4.add_paragraph()
p2.text = "2. Automated error fixes (no bureaucratic delays or ticket queues)"
p2.font.size = Pt(16)
p2.font.color.rgb = TEXT_WHITE

p3 = tf4.add_paragraph()
p3.text = "3. Targeted, agile marketing & rapid execution of new ideas"
p3.font.size = Pt(16)
p3.font.color.rgb = TEXT_WHITE

p4 = tf4.add_paragraph()
p4.text = "4. Maximum customer customization (100% tailored vs rigid templates)"
p4.font.size = Pt(16)
p4.font.color.rgb = TEXT_WHITE

# Slide 5: Personal Commitment
slide5 = prs.slides.add_slide(blank_layout)
add_header(slide5, "Slide 05 / Leadership Commitment", "True Loyalty & Skin in the Game", "A Captain Stays With His Ship")

tb5 = slide5.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(4.0))
tf5 = tb5.text_frame
tf5.word_wrap = True

p = tf5.paragraphs[0]
p.text = "\"A captain stays with his ship. My commitment to this vision and to our people has always been 100% absolute. Brand is important, but PEOPLE are paramount.\""
p.font.size = Pt(20)
p.font.italic = True
p.font.color.rgb = CYAN_GLOW

p2 = tf5.add_paragraph()
p2.text = "• Shared Lifeline: This business is my primary focus and lifeline too—we rise and win together."
p2.font.size = Pt(16)
p2.font.color.rgb = TEXT_WHITE

# Slide 6: DP World Reality
slide6 = prs.slides.add_slide(blank_layout)
add_header(slide6, "Slide 06 / Crucial Context", "Taking Ownership of Our Future", "The DP World Reality: Our Unstoppable Catalyst")

tb6 = slide6.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(4.0))
tf6 = tb6.text_frame
tf6.word_wrap = True

p = tf6.paragraphs[0]
p.text = "• Official Corporate Reality: DP World is shutting down service extensions and client renewals."
p.font.size = Pt(18)
p.font.color.rgb = AMBER

p2 = tf6.add_paragraph()
p2.text = "• Formal Steps Underway: Confirming notice terms to notify clients transparently."
p2.font.size = Pt(18)
p2.font.color.rgb = TEXT_WHITE

p3 = tf6.add_paragraph()
p3.text = "• The Unstoppable Takeaway: We no longer depend on external corporate decisions. We own our destiny right now."
p3.font.size = Pt(18)
p3.font.color.rgb = EMERALD

# Slide 7: The 2-Month Sprint
slide7 = prs.slides.add_slide(blank_layout)
add_header(slide7, "Slide 07 / Mobilization", "Time to Mobilize: The 60-Day Push", "Bite the Bullet: High Intensity Sprint Mode")

tb7 = slide7.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(4.0))
tf7 = tb7.text_frame
tf7.word_wrap = True

p = tf7.paragraphs[0]
p.text = "• Intense Mobilization: The next 2 months require going beyond standard 8-hour workdays. Every hour counts."
p.font.size = Pt(18)
p.font.color.rgb = AMBER

p2 = tf7.add_paragraph()
p2.text = "• Mandatory Sprint: Essential to put us firmly on our feet and establish unstoppable market momentum."
p2.font.size = Pt(18)
p2.font.color.rgb = TEXT_WHITE

# Slide 8: Three Pillars
slide8 = prs.slides.add_slide(blank_layout)
add_header(slide8, "Slide 08 / Execution Standards", "The Three Pillars of Execution", "Our Non-Negotiable Standards for the Sprint")

tb8 = slide8.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(4.0))
tf8 = tb8.text_frame
tf8.word_wrap = True

p = tf8.paragraphs[0]
p.text = "1. DISCIPLINE: Zero excuses, no 'I forgot' moments, total operational precision and punctuality."
p.font.size = Pt(18)
p.font.color.rgb = TEXT_WHITE

p2 = tf8.add_paragraph()
p2.text = "2. RESPONSIBILITY: End-to-end ownership of every lead, client, and issue until 100% resolved."
p2.font.size = Pt(18)
p2.font.color.rgb = TEXT_WHITE

p3 = tf8.add_paragraph()
p3.text = "3. MARKETING INTEGRATION: Aggressive outreach, loud market presence, and precise targeted growth."
p3.font.size = Pt(18)
p3.font.color.rgb = TEXT_WHITE

# Slide 9: Spartan Mode
slide9 = prs.slides.add_slide(blank_layout)
add_header(slide9, "Slide 09 / Culture", "Spartan Mode Activation", "High Intensity, Zero Friction, Total Accountability")

tb9 = slide9.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(4.0))
tf9 = tb9.text_frame
tf9.word_wrap = True

p = tf9.paragraphs[0]
p.text = "\"We embrace a spartan, high-performance culture where speed, grit, and accountability drive daily wins. Zero drama, zero friction, maximum focus.\""
p.font.size = Pt(20)
p.font.italic = True
p.font.color.rgb = AMBER

# Slide 10: Closing
slide10 = prs.slides.add_slide(blank_layout)
set_slide_background(slide10)

tb10 = slide10.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.0))
tf10 = tb10.text_frame
tf10.word_wrap = True

p_tag = tf10.paragraphs[0]
p_tag.text = "SLIDE 10 / CLOSING"
p_tag.font.size = Pt(14)
p_tag.font.bold = True
p_tag.font.color.rgb = CYAN_GLOW

p_title = tf10.add_paragraph()
p_title.text = "We Have Everything We Need to Win"
p_title.font.size = Pt(36)
p_title.font.bold = True
p_title.font.color.rgb = TEXT_WHITE

p_sub = tf10.add_paragraph()
p_sub.text = "This is not just about company survival—it's about personal growth, financial mastery, and reaching a whole new level together."
p_sub.font.size = Pt(18)
p_sub.font.color.rgb = TEXT_MUTED

p_call = tf10.add_paragraph()
p_call.text = "Let's build this future NOW."
p_call.font.size = Pt(24)
p_call.font.bold = True
p_call.font.color.rgb = CYAN_GLOW

prs.save(pptx_path)
print(f"Successfully generated PPTX presentation at: {pptx_path}")
