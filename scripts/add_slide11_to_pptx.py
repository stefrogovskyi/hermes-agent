import pptx, os
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

pptx_path = r'C:\Users\Stefan\AppData\Local\hermes\reports\navo24_team_presentation.pptx'

prs = pptx.Presentation(pptx_path)
blank_layout = prs.slide_layouts[6]

BG_COLOR = RGBColor(11, 15, 25)
CYAN_GLOW = RGBColor(56, 189, 248)
TEXT_WHITE = RGBColor(248, 250, 252)
TEXT_MUTED = RGBColor(148, 163, 184)
EMERALD = RGBColor(16, 185, 129)

slide11 = prs.slides.add_slide(blank_layout)
slide11.background.fill.solid()
slide11.background.fill.fore_color.rgb = BG_COLOR

# Header
tb = slide11.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.8))
tf = tb.text_frame
tf.word_wrap = True

p_tag = tf.paragraphs[0]
p_tag.text = "SLIDE 11 / FINANCIAL AUDIT & BASELINE"
p_tag.font.size = Pt(12)
p_tag.font.bold = True
p_tag.font.color.rgb = CYAN_GLOW

p_title = tf.add_paragraph()
p_title.text = "Current Sales Revenue & Team Portfolio"
p_title.font.size = Pt(28)
p_title.font.bold = True
p_title.font.color.rgb = TEXT_WHITE

p_sub = tf.add_paragraph()
p_sub.text = "A Solid $345,000+ Monthly Revenue Foundation Managed by 8 Consultants"
p_sub.font.size = Pt(14)
p_sub.font.color.rgb = TEXT_MUTED

# Content
tb_c = slide11.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.7), Inches(4.5))
tf_c = tb_c.text_frame
tf_c.word_wrap = True

lines = [
    "• Monthly Cash-Flow (PAYG + Monthly): $167,122 / mo (Factual active monthly recurring stream)",
    "• Annual Contract Portfolio (ARR): $2,034,223 / yr ($198,685 / mo amortized equivalent)",
    "• Total Combined Gross Stream: $345,807 / mo equivalent across 8 B2B Sales Consultants",
    "",
    "Consultants Breakdown:",
    "  - Katya Komarova: $50,000 / mo | $600,000 ARR (Contract to Sept 2026/2027)",
    "  - Sasha Grabarchuk: $58,844 / mo | $700,000 ARR ($38.9k/mo due Oct 2026)",
    "  - Lilya Khovrak: $30,000 / mo | $150,000 ARR (Strong recurring base)",
    "  - Lera Guliy: $13,500 / mo | $300,000 ARR ($300k ARR due Oct 3 2026)",
    "  - Andrey Gorodinsky: $11,000 / mo | $92,000 ARR (Nauta, First, Herpot $15k/qtr Sept 2026)",
    "  - Oleg Chervinsky: $1,478 / mo | $39,433 ARR (Tradewind $4.4k trial ends Aug 2026)",
    "  - Katya Kernesh: $1,300 / mo | $100,790 ARR ($57.6k ARR due Nov 14 2026)",
    "  - Katya Kapustyan: $1,000 / mo | $52,000 ARR ($52k ARR due April/June 2027)"
]

for idx, line in enumerate(lines):
    p = tf_c.paragraphs[0] if idx == 0 else tf_c.add_paragraph()
    p.text = line
    p.font.size = Pt(13) if "Consultants Breakdown" not in line else Pt(14)
    if "Monthly Cash-Flow" in line or "Total Combined" in line:
        p.font.bold = True
        p.font.color.rgb = EMERALD
    elif "Consultants Breakdown" in line:
        p.font.bold = True
        p.font.color.rgb = CYAN_GLOW
    else:
        p.font.color.rgb = TEXT_WHITE

prs.save(pptx_path)
print('✅ Successfully added Slide 11 to navo24_team_presentation.pptx!')
